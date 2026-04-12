from __future__ import annotations

import io
import json
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from docx import Document
from matplotlib import pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from ml_service.core.config import Settings


def ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def format_timestamp_srt(milliseconds: int) -> str:
    total_seconds, ms = divmod(max(milliseconds, 0), 1000)
    minutes, seconds = divmod(total_seconds, 60)
    hours, minutes = divmod(minutes, 60)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d},{ms:03d}"


def build_srt(segments: list[dict[str, Any]]) -> str:
    lines: list[str] = []
    for index, segment in enumerate(segments, start=1):
        start_ms = int(segment.get("start_ms") or 0)
        end_ms = int(segment.get("end_ms") or max(start_ms + 1000, start_ms))
        lines.extend(
            [
                str(index),
                f"{format_timestamp_srt(start_ms)} --> {format_timestamp_srt(end_ms)}",
                str(segment.get("text") or ""),
                "",
            ]
        )
    return "\n".join(lines).strip() + "\n"


def probe_media_duration_seconds(settings: Settings, source_path: Path) -> float | None:
    command = [
        settings.ffprobe_binary,
        "-v",
        "error",
        "-show_entries",
        "format=duration",
        "-of",
        "default=noprint_wrappers=1:nokey=1",
        str(source_path),
    ]
    try:
        result = subprocess.run(command, check=True, capture_output=True, text=True)
    except Exception:  # noqa: BLE001
        return None
    output = result.stdout.strip()
    if not output:
        return None
    try:
        duration = float(output)
    except ValueError:
        return None
    return duration if duration >= 0 else None


def create_text_pdf(title: str, sections: list[tuple[str, str]]) -> bytes:
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4
    cursor = height - 60
    pdf.setFont("Helvetica-Bold", 18)
    pdf.drawString(50, cursor, title[:90])
    cursor -= 30
    pdf.setFont("Helvetica", 11)
    for heading, body in sections:
        if cursor < 90:
            pdf.showPage()
            cursor = height - 60
            pdf.setFont("Helvetica", 11)
        pdf.setFont("Helvetica-Bold", 13)
        pdf.drawString(50, cursor, heading[:90])
        cursor -= 20
        pdf.setFont("Helvetica", 11)
        for line in wrap_text(body, 92):
            if cursor < 70:
                pdf.showPage()
                cursor = height - 60
                pdf.setFont("Helvetica", 11)
            pdf.drawString(50, cursor, line)
            cursor -= 16
        cursor -= 10
    pdf.save()
    return buffer.getvalue()


def create_docx(title: str, sections: list[tuple[str, str]]) -> bytes:
    document = Document()
    document.add_heading(title, 0)
    for heading, body in sections:
        document.add_heading(heading, level=1)
        document.add_paragraph(body)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def presentation_to_pptx(project: dict[str, Any], output_path: Path) -> None:
    ensure_parent(output_path)
    presentation = Presentation()
    for slide_data in project.get("slides", []):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(slide_data.get("title") or "Slide")
        subtitle = slide.placeholders[1]
        blocks = slide_data.get("blocks") or []
        bullet_items: list[str] = []
        for block in blocks:
            data = block.get("data") or {}
            if block.get("kind") == "bullet_list":
                bullet_items.extend([str(item) for item in data.get("items", [])])
            elif block.get("kind") == "text":
                bullet_items.append(str(data.get("content") or ""))
        subtitle.text = "\n".join(item for item in bullet_items if item)[:2000]
        if slide_data.get("speaker_notes"):
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = str(slide_data.get("speaker_notes"))
    presentation.save(str(output_path))


def presentation_to_html(project: dict[str, Any]) -> str:
    title = project.get("title") or "Presentation"
    slides_html = []
    for slide in project.get("slides", []):
        blocks_html = []
        for block in slide.get("blocks", []):
            kind = block.get("kind")
            data = block.get("data") or {}
            if kind == "bullet_list":
                items = "".join(f"<li>{item}</li>" for item in data.get("items", []))
                blocks_html.append(f"<ul>{items}</ul>")
            elif kind == "text":
                blocks_html.append(f"<p>{data.get('content', '')}</p>")
            else:
                blocks_html.append(f"<pre>{json.dumps(data, ensure_ascii=True)}</pre>")
        slides_html.append(
            "<section class='slide'>"
            f"<h2>{slide.get('title', '')}</h2>"
            f"<h3>{slide.get('subtitle', '')}</h3>"
            + "".join(blocks_html)
            + f"<aside>{slide.get('speaker_notes', '')}</aside>"
            + "</section>"
        )
    return (
        "<!DOCTYPE html><html><head><meta charset='utf-8'><title>"
        + str(title)
        + "</title><style>body{font-family:Arial,sans-serif;background:#f8fafc;color:#0f172a;padding:40px;}section.slide{background:#fff;border:1px solid #cbd5e1;border-radius:16px;padding:24px;margin:0 0 24px;}aside{margin-top:16px;color:#475569;white-space:pre-wrap;}ul{padding-left:20px;}</style></head><body>"
        + f"<h1>{title}</h1>"
        + "".join(slides_html)
        + "</body></html>"
    )


def render_slide_image(slide: dict[str, Any], output_path: Path, size: tuple[int, int] = (1600, 900)) -> None:
    ensure_parent(output_path)
    image = Image.new("RGB", size, color=(248, 250, 252))
    draw = ImageDraw.Draw(image)
    title_font = ImageFont.load_default()
    body_font = ImageFont.load_default()
    draw.rectangle((40, 40, size[0] - 40, size[1] - 40), outline=(14, 116, 144), width=4)
    draw.text((80, 80), str(slide.get("title") or "Slide"), fill=(15, 23, 42), font=title_font)
    y = 150
    for block in slide.get("blocks", []):
        data = block.get("data") or {}
        if block.get("kind") == "bullet_list":
            for item in data.get("items", []):
                for line in wrap_text(str(item), 70):
                    draw.text((100, y), f"- {line}", fill=(30, 41, 59), font=body_font)
                    y += 28
        elif block.get("kind") == "text":
            for line in wrap_text(str(data.get("content") or ""), 80):
                draw.text((100, y), line, fill=(51, 65, 85), font=body_font)
                y += 24
        y += 10
    image.save(output_path)


def wrap_text(text: str, width: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        if len(current) + 1 + len(word) <= width:
            current = f"{current} {word}"
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def convert_with_libreoffice(settings: Settings, input_path: Path, target_ext: str) -> bytes | None:
    output_dir = Path(tempfile.mkdtemp())
    command = [
        settings.libreoffice_binary,
        "--headless",
        "--convert-to",
        target_ext.lstrip("."),
        str(input_path),
        "--outdir",
        str(output_dir),
    ]
    try:
        subprocess.run(command, check=True, capture_output=True)
    except Exception:  # noqa: BLE001
        return None
    converted = output_dir / f"{input_path.stem}.{target_ext.lstrip('.')}"
    if not converted.exists():
        return None
    return converted.read_bytes()


def dataframe_to_csv_bytes(columns: list[dict[str, Any]], rows: list[dict[str, Any]]) -> bytes:
    keys = [column.get("key") or column.get("label") for column in columns]
    frame = pd.DataFrame(rows, columns=keys)
    return frame.to_csv(index=False).encode("utf-8")


def dataframe_to_xlsx_bytes(columns: list[dict[str, Any]], rows: list[dict[str, Any]]) -> bytes:
    keys = [column.get("key") or column.get("label") for column in columns]
    frame = pd.DataFrame(rows, columns=keys)
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        frame.to_excel(writer, index=False, sheet_name="data")
    return buffer.getvalue()


def infographic_to_image_bytes(spec: dict[str, Any], output_format: str) -> bytes:
    figure, axis = plt.subplots(figsize=(10, 6))
    axis.set_title(str(spec.get("title") or "Infographic"))
    blocks = spec.get("blocks") or []
    labels = [str(block.get("label") or block.get("kind")) for block in blocks]
    values = [float(block.get("value") or index + 1) for index, block in enumerate(blocks)]
    if not labels:
        labels = ["placeholder"]
        values = [1]
    axis.bar(labels, values, color=["#0F766E", "#F59E0B", "#1D4ED8", "#DC2626", "#6D28D9"][: len(labels)])
    axis.tick_params(axis="x", rotation=25)
    axis.set_ylabel("Value")
    buffer = io.BytesIO()
    figure.tight_layout()
    figure.savefig(buffer, format=output_format)
    plt.close(figure)
    return buffer.getvalue()


def ffmpeg_concat_audio(settings: Settings, inputs: list[Path], output_path: Path) -> None:
    ensure_parent(output_path)
    manifest = output_path.with_suffix(".txt")
    manifest.write_text("\n".join(f"file '{path.as_posix()}'" for path in inputs), encoding="utf-8")
    command = [
        settings.ffmpeg_binary,
        "-y",
        "-f",
        "concat",
        "-safe",
        "0",
        "-i",
        str(manifest),
        "-af",
        "loudnorm",
        str(output_path),
    ]
    subprocess.run(command, check=True, capture_output=True)
    manifest.unlink(missing_ok=True)


def ffmpeg_convert(settings: Settings, input_path: Path, output_path: Path) -> None:
    ensure_parent(output_path)
    command = [settings.ffmpeg_binary, "-y", "-i", str(input_path), str(output_path)]
    subprocess.run(command, check=True, capture_output=True)


def ffmpeg_render_audiogram(settings: Settings, audio_path: Path, output_path: Path, cover_path: Path | None = None) -> None:
    ensure_parent(output_path)
    if cover_path and cover_path.exists():
        command = [
            settings.ffmpeg_binary,
            "-y",
            "-loop",
            "1",
            "-i",
            str(cover_path),
            "-i",
            str(audio_path),
            "-filter_complex",
            "[1:a]showwaves=s=1280x720:mode=line:colors=0F766E,format=rgba[sw];[0:v][sw]overlay=shortest=1",
            "-c:v",
            "libx264",
            "-c:a",
            "aac",
            "-shortest",
            str(output_path),
        ]
    else:
        command = [
            settings.ffmpeg_binary,
            "-y",
            "-i",
            str(audio_path),
            "-filter_complex",
            "showwaves=s=1280x720:mode=line:colors=0F766E",
            "-c:v",
            "libx264",
            "-c:a",
            "aac",
            "-shortest",
            str(output_path),
        ]
    subprocess.run(command, check=True, capture_output=True)


def ffmpeg_render_slide_segment(settings: Settings, image_path: Path, audio_path: Path, output_path: Path) -> None:
    ensure_parent(output_path)
    command = [
        settings.ffmpeg_binary,
        "-y",
        "-loop",
        "1",
        "-i",
        str(image_path),
        "-i",
        str(audio_path),
        "-c:v",
        "libx264",
        "-tune",
        "stillimage",
        "-c:a",
        "aac",
        "-b:a",
        "192k",
        "-pix_fmt",
        "yuv420p",
        "-shortest",
        str(output_path),
    ]
    subprocess.run(command, check=True, capture_output=True)


def ffmpeg_concat_video(settings: Settings, inputs: list[Path], output_path: Path) -> None:
    ensure_parent(output_path)
    manifest = output_path.with_suffix(".txt")
    manifest.write_text("\n".join(f"file '{path.as_posix()}'" for path in inputs), encoding="utf-8")
    command = [
        settings.ffmpeg_binary,
        "-y",
        "-f",
        "concat",
        "-safe",
        "0",
        "-i",
        str(manifest),
        "-c",
        "copy",
        str(output_path),
    ]
    subprocess.run(command, check=True, capture_output=True)
    manifest.unlink(missing_ok=True)


def create_cover_image(title: str, output_path: Path, size: tuple[int, int] = (1280, 720)) -> None:
    ensure_parent(output_path)
    image = Image.new("RGB", size, color=(15, 23, 42))
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()
    y = size[1] // 3
    for line in wrap_text(title, 32):
        draw.text((80, y), line, fill=(248, 250, 252), font=font)
        y += 28
    image.save(output_path)
